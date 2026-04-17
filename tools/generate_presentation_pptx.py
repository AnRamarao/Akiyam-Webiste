#!/usr/bin/env python3
"""
Generate AIKYAM Architecture Presentation (PowerPoint .pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
NAVY = RGBColor(0x1A, 0x36, 0x5D)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
AMBER = RGBColor(0xFA, 0xCC, 0x15)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_text(slide, left, top, width, height, text, font_size=12,
                        fill_color=CARD_BG, font_color=WHITE, bold=False,
                        alignment=PP_ALIGN.LEFT, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    if isinstance(text, str):
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
    elif isinstance(text, list):
        for i, line in enumerate(text):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            if isinstance(line, tuple):
                p.text = line[0]
                p.font.size = Pt(line[1]) if len(line) > 1 else Pt(font_size)
                p.font.color.rgb = line[2] if len(line) > 2 else font_color
                p.font.bold = line[3] if len(line) > 3 else False
            else:
                p.text = line
                p.font.size = Pt(font_size)
                p.font.color.rgb = font_color
            p.alignment = alignment

    return shape

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Gold accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(18)
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.LEFT

    return slide

def add_section_slide(prs, title_text, subtitle_text=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = GOLD
        p2.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title_text, slide_num=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top gold line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.04), Inches(10), Inches(0.75))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0x12, 0x1B, 0x30)
    title_bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(8.8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.color.rgb = GOLD
    p.font.bold = True

    return slide

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    num_rows = 1 + len(rows)
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = WHITE
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x16, 0x20, 0x33)

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    return table

def add_bullet_list(slide, left, top, width, height, items, font_size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            p.text = item[0]
            p.font.size = Pt(item[1]) if len(item) > 1 else Pt(font_size)
            p.font.color.rgb = item[2] if len(item) > 2 else WHITE
            p.font.bold = item[3] if len(item) > 3 else False
        else:
            p.text = f"\u2022  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
        p.space_after = Pt(6)

    return txBox


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ===== SLIDE 1: TITLE =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Gold line top
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'AIKYAM'
    p.font.size = Pt(54)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = 'Technical Architecture & Event Management'
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT

    # Meta info
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.0))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = 'April 2026  |  Ramarao  |  Team Presentation'
    p3.font.size = Pt(14)
    p3.font.color.rgb = MED_GRAY
    p3.alignment = PP_ALIGN.LEFT

    # ===== SLIDE 2: AGENDA =====
    slide = add_content_slide(prs, 'Agenda')
    items = [
        'Platform Overview',
        'Architecture at a Glance',
        'Technology Stack',
        'Authentication System',
        'Role-Based Access Control (9 Levels)',
        'Event Management Feature',
        'Database Schema & Security',
        'Developer Workflow & Demo',
        'What\'s Next',
    ]
    for i, item in enumerate(items):
        y = Inches(1.1) + Inches(i * 0.45)
        # Number circle
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.8), y, Inches(0.35), Inches(0.35))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = GOLD
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK_BG
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Text
        txBox = slide.shapes.add_textbox(Inches(1.4), y, Inches(7), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    # ===== SLIDE 3: PLATFORM OVERVIEW =====
    slide = add_content_slide(prs, 'Platform Overview')

    add_shape_with_text(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(2.2),
        [('What We Manage', 14, GOLD, True),
         ('', 6, WHITE, False),
         ('\u2022  Community Events', 12, WHITE, False),
         ('\u2022  Team & Board Info', 12, WHITE, False),
         ('\u2022  Photo Gallery', 12, WHITE, False),
         ('\u2022  Vendor Directory', 12, WHITE, False),
         ('\u2022  Member Authentication', 12, WHITE, False),
         ('\u2022  Donations', 12, WHITE, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    add_shape_with_text(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(2.2),
        [('How We Build It', 14, GOLD, True),
         ('', 6, WHITE, False),
         ('\u2022  HTML / CSS / Vanilla JS', 12, WHITE, False),
         ('\u2022  No frameworks or build tools', 12, WHITE, False),
         ('\u2022  GoDaddy static hosting', 12, WHITE, False),
         ('\u2022  Supabase free tier backend', 12, WHITE, False),
         ('\u2022  Zero monthly cost', 12, WHITE, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # Key principle box
    add_shape_with_text(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.8),
        'KEY PRINCIPLE:  Zero-Cost Infrastructure \u2014 No servers, no monthly bills, no framework lock-in',
        font_size=13, fill_color=RGBColor(0x33, 0x2B, 0x00), font_color=GOLD,
        bold=True, alignment=PP_ALIGN.CENTER, border_color=GOLD)

    # ===== SLIDE 4: ARCHITECTURE DIAGRAM =====
    slide = add_content_slide(prs, 'Architecture at a Glance')

    # Browser box
    add_shape_with_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(1.8),
        [("USER'S BROWSER", 11, GOLD, True),
         ('', 4, WHITE, False),
         ('index.html  |  styles.css  |  script.js  |  auth.js  |  events.js  |  admin-events.js', 10, WHITE, False),
         ('', 4, WHITE, False),
         ('Everything runs client-side \u2014 no backend server needed', 10, MED_GRAY, False)],
        alignment=PP_ALIGN.CENTER, border_color=GOLD)

    # Arrow
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.95), Inches(1), Inches(0.3))
    txBox.text_frame.paragraphs[0].text = '\u25bc'
    txBox.text_frame.paragraphs[0].font.color.rgb = GOLD
    txBox.text_frame.paragraphs[0].font.size = Pt(18)

    txBox2 = slide.shapes.add_textbox(Inches(7), Inches(2.95), Inches(1), Inches(0.3))
    txBox2.text_frame.paragraphs[0].text = '\u25bc'
    txBox2.text_frame.paragraphs[0].font.color.rgb = GOLD
    txBox2.text_frame.paragraphs[0].font.size = Pt(18)

    # GoDaddy
    add_shape_with_text(slide, Inches(0.5), Inches(3.3), Inches(4.2), Inches(1.5),
        [('GoDaddy  (Static Host)', 12, WHITE, True),
         ('', 4, WHITE, False),
         ('HTML, CSS, JS, Images, JSON data', 10, LIGHT_GRAY, False),
         ('No server-side processing', 10, MED_GRAY, False)],
        alignment=PP_ALIGN.CENTER, fill_color=RGBColor(0x1A, 0x36, 0x5D),
        border_color=RGBColor(0x33, 0x55, 0x88))

    # Supabase
    add_shape_with_text(slide, Inches(5.3), Inches(3.3), Inches(4.2), Inches(1.5),
        [('Supabase  (Free Tier)', 12, WHITE, True),
         ('', 4, WHITE, False),
         ('Auth (GoTrue)  |  PostgreSQL  |  RLS', 10, LIGHT_GRAY, False),
         ('Server-side authorization enforcement', 10, MED_GRAY, False)],
        alignment=PP_ALIGN.CENTER, fill_color=RGBColor(0x1A, 0x36, 0x5D),
        border_color=RGBColor(0x3E, 0xCF, 0x8E))

    # ===== SLIDE 5: TECH STACK =====
    slide = add_content_slide(prs, 'Technology Stack')
    add_table(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(3.5),
        ['Layer', 'Technology', 'Details'],
        [
            ['Markup', 'HTML5', 'Semantic elements, ARIA attributes'],
            ['Styling', 'CSS3', 'Custom properties, Grid, Flexbox, BEM naming'],
            ['Logic', 'Vanilla JavaScript', 'ES2017+ (async/await), IIFE modules'],
            ['Fonts', 'Google Fonts', 'Montserrat (400, 700, 900)'],
            ['Auth', 'Supabase Auth', 'Client-side SDK v2 via CDN'],
            ['Database', 'PostgreSQL', 'Supabase free tier, Row-Level Security'],
            ['Hosting', 'GoDaddy', 'Static file hosting (included with domain)'],
            ['Build', 'Make + Python3', 'Image optimization, CSS/JS minification'],
        ],
        col_widths=[1.5, 2.5, 5.0])

    # ===== SLIDE 6: MODULE ARCHITECTURE =====
    slide = add_content_slide(prs, 'JavaScript Module Architecture')

    modules = [
        (Inches(3.8), Inches(1.1), 'config.js', 'Credentials', MED_GRAY),
        (Inches(3.8), Inches(1.8), 'auth.js', 'AIKYAM_Auth', TEAL),
        (Inches(3.8), Inches(2.6), 'events.js', 'AIKYAM_Events', GREEN),
        (Inches(1.2), Inches(3.5), 'script.js', 'Public site rendering', GOLD),
        (Inches(5.8), Inches(3.5), 'admin-events.js', 'Admin dashboard', RED),
    ]

    for x, y, name, desc, color in modules:
        add_shape_with_text(slide, x, y, Inches(2.4), Inches(0.55),
            [( name, 11, color, True),
             (desc, 8, MED_GRAY, False)],
            alignment=PP_ALIGN.CENTER,
            border_color=color)

    # Arrows
    for arrow_y in [Inches(1.65), Inches(2.4), Inches(3.15)]:
        txBox = slide.shapes.add_textbox(Inches(4.7), arrow_y, Inches(0.5), Inches(0.25))
        txBox.text_frame.paragraphs[0].text = '\u25bc'
        txBox.text_frame.paragraphs[0].font.color.rgb = GOLD
        txBox.text_frame.paragraphs[0].font.size = Pt(12)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Load order note
    add_shape_with_text(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.5),
        'Load Order (all defer):  supabase-js (CDN)  \u2192  config.js  \u2192  auth.js  \u2192  events.js  \u2192  script.js',
        font_size=10, fill_color=RGBColor(0x16, 0x20, 0x33), font_color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 7: AUTH SYSTEM =====
    slide = add_content_slide(prs, 'Authentication System')

    # Signup flow
    add_shape_with_text(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(3.5),
        [('Sign-Up Flow', 13, GOLD, True),
         ('', 4, WHITE, False),
         ('1. User fills signup form', 10, WHITE, False),
         ('2. Rate limit check (5/60s)', 10, WHITE, False),
         ('3. Supabase creates account', 10, WHITE, False),
         ('4. Confirmation email sent', 10, WHITE, False),
         ('5. User clicks email link', 10, WHITE, False),
         ('6. auth.js detects hash fragment', 10, WHITE, False),
         ('7. Profile auto-created (Member)', 10, WHITE, False),
         ('8. User can now log in', 10, WHITE, False),
         ('', 6, WHITE, False),
         ('Security Features:', 11, GOLD, True),
         ('\u2022  Bcrypt password hashing', 9, LIGHT_GRAY, False),
         ('\u2022  Rate limiting (5/min)', 9, LIGHT_GRAY, False),
         ('\u2022  Email confirmation required', 9, LIGHT_GRAY, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # Login flow
    add_shape_with_text(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(3.5),
        [('Login Flow', 13, GOLD, True),
         ('', 4, WHITE, False),
         ('1. Email + Password submitted', 10, WHITE, False),
         ('2. Supabase validates credentials', 10, WHITE, False),
         ('3. Session token returned', 10, WHITE, False),
         ('4. Header UI updates:', 10, WHITE, False),
         ('   Login btn \u2192 Avatar dropdown', 10, TEAL, False),
         ('5. Role-gated UI appears', 10, WHITE, False),
         ('', 6, WHITE, False),
         ('Route Protection:', 11, GOLD, True),
         ('\u2022  No session \u2192 redirect to login', 9, LIGHT_GRAY, False),
         ('\u2022  Insufficient role \u2192 redirect home', 9, LIGHT_GRAY, False),
         ('\u2022  Valid role \u2192 page loads', 9, LIGHT_GRAY, False),
         ('', 6, WHITE, False),
         ('data-auth-required="ProgramManager"', 8, MED_GRAY, False),
         ('hides/shows elements per role', 8, MED_GRAY, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # ===== SLIDE 8: ROLE HIERARCHY =====
    slide = add_content_slide(prs, 'Role-Based Access Control \u2014 9-Level Hierarchy')

    roles_data = [
        ('8', 'GlobalAdmin', 'Full system access, manage all users'),
        ('7', 'OrgAdmin', 'Org-level admin, manage team & events'),
        ('6', 'ContentManager', 'Content lifecycle (publish/deactivate)'),
        ('5', 'ProgramManager', 'Event creation and management'),
        ('4', 'FinanceStaff', 'Financial access and reporting'),
        ('3', 'Volunteer', 'Event coordination, member access'),
        ('2', 'Donor', 'Donation history, donor content'),
        ('1', 'Support', 'Support queue access'),
        ('0', 'Member', 'Default role, basic access'),
    ]

    for i, (level, role, desc) in enumerate(roles_data):
        y = Inches(1.05) + Inches(i * 0.44)
        # Width shrinks as level decreases (visual hierarchy)
        w = Inches(5.5 + (int(level) * 0.3))
        bar_color = RGBColor(
            max(0x1A, min(0xFF, 0x1A + int(level) * 28)),
            max(0x36, min(0xD7, 0x36 + int(level) * 18)),
            max(0x5D, min(0xFF, 0x5D - int(level) * 5))
        )

        add_shape_with_text(slide, Inches(0.5), y, w, Inches(0.38),
            f'  Level {level}  |  {role}  \u2014  {desc}',
            font_size=9, fill_color=bar_color, font_color=WHITE, bold=False)

    # How it works note
    add_shape_with_text(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4),
        'hasMinRole(\'ContentManager\', \'ProgramManager\') \u2192 true (6 >= 5)  |  Enforced client-side AND server-side via RLS',
        font_size=9, fill_color=RGBColor(0x16, 0x20, 0x33), font_color=GOLD,
        alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 9: EVENT LIFECYCLE =====
    slide = add_content_slide(prs, 'Event Management \u2014 Lifecycle')

    # State boxes
    states = [
        (Inches(0.8), Inches(1.5), 'DRAFT', 'Created by ProgramManager+\nDefault state on creation', MED_GRAY),
        (Inches(3.5), Inches(1.5), 'PUBLISHED', 'Visible on public site\nManual or scheduled publish', GREEN),
        (Inches(6.5), Inches(1.5), 'EXPIRED', 'Auto when end_at passes\nShown as past events', MED_GRAY),
        (Inches(3.5), Inches(3.3), 'DEACTIVATED', 'Manually hidden\nCan be re-published', AMBER),
    ]

    for x, y, state, desc, color in states:
        add_shape_with_text(slide, x, y, Inches(2.5), Inches(1.2),
            [(state, 14, color, True),
             ('', 2, WHITE, False),
             (desc, 9, LIGHT_GRAY, False)],
            alignment=PP_ALIGN.CENTER, border_color=color)

    # Arrows between states
    arrows = [
        (Inches(3.1), Inches(1.9), '\u25b6', 'Publish'),
        (Inches(5.85), Inches(1.9), '\u25b6', 'Auto-expire'),
        (Inches(4.5), Inches(2.75), '\u25bc', 'Deactivate'),
    ]
    for x, y, symbol, label in arrows:
        txBox = slide.shapes.add_textbox(x, y, Inches(0.7), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f'{symbol} {label}'
        p.font.size = Pt(8)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

    # Auto-lifecycle note
    add_shape_with_text(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(0.5),
        'Auto-Lifecycle:  auto_expire_events() + auto_publish_events() run every 15 min via pg_cron + on every page load',
        font_size=9, fill_color=RGBColor(0x16, 0x20, 0x33), font_color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 10: EVENT PERMISSIONS =====
    slide = add_content_slide(prs, 'Event Management \u2014 Permissions')

    add_table(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(2.5),
        ['Role', 'View All', 'Create/Edit', 'Publish/Deactivate', 'Delete'],
        [
            ['GlobalAdmin', 'Yes', 'Yes', 'Yes', 'Yes'],
            ['OrgAdmin', 'Yes', 'Yes', 'Yes', 'Yes'],
            ['ContentManager', 'Yes', 'Yes', 'Yes', 'No'],
            ['ProgramManager', 'Yes', 'Own only', 'No', 'No'],
            ['Below PM', 'No access to admin page', '', '', ''],
        ],
        col_widths=[2.0, 1.5, 2.0, 2.0, 1.5])

    # Fallback box
    add_shape_with_text(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(1.2),
        [('Data Fallback Strategy', 12, GOLD, True),
         ('', 4, WHITE, False),
         ('Primary: Supabase events table (live data, real-time lifecycle)', 10, WHITE, False),
         ('Fallback: Static JSON files (data/upcomingEvents.json + completedEvents.json)', 10, WHITE, False),
         ('The site never shows a blank page \u2014 same rendering code works with either source', 10, MED_GRAY, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # ===== SLIDE 11: DATABASE SCHEMA =====
    slide = add_content_slide(prs, 'Database Schema')

    # Profiles box
    add_shape_with_text(slide, Inches(0.5), Inches(1.1), Inches(3.0), Inches(2.0),
        [('profiles', 13, TEAL, True),
         ('', 2, WHITE, False),
         ('id (UUID, PK, FK \u2192 auth.users)', 8, LIGHT_GRAY, False),
         ('email (TEXT)', 8, LIGHT_GRAY, False),
         ('full_name (TEXT)', 8, LIGHT_GRAY, False),
         ('role (app_role ENUM)', 8, GOLD, False),
         ('created_at, updated_at', 8, MED_GRAY, False),
         ('', 2, WHITE, False),
         ('Trigger: auto-create on signup', 7, MED_GRAY, False)],
        border_color=TEAL)

    # Events box
    add_shape_with_text(slide, Inches(3.8), Inches(1.1), Inches(5.8), Inches(2.0),
        [('events', 13, GREEN, True),
         ('', 2, WHITE, False),
         ('id, title, description, location, price, img, tbd', 8, LIGHT_GRAY, False),
         ('start_at, end_at (TIMESTAMPTZ)', 8, LIGHT_GRAY, False),
         ('status (event_status: draft|published|deactivated|expired)', 8, GOLD, False),
         ('publish_at, expired_at, summary', 8, LIGHT_GRAY, False),
         ('created_by, updated_by (FK \u2192 profiles)', 8, LIGHT_GRAY, False),
         ('', 2, WHITE, False),
         ('Triggers: auto-update updated_at  |  Indexes: status+start_at, publish_at', 7, MED_GRAY, False)],
        border_color=GREEN)

    # RLS summary
    add_shape_with_text(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(1.8),
        [('Row-Level Security (RLS) \u2014 Server-Side Authorization', 11, GOLD, True),
         ('', 3, WHITE, False),
         ('Profiles: Users read own | Admins read all | Staff read all | Users update own name | GlobalAdmin updates any role', 8, LIGHT_GRAY, False),
         ('', 2, WHITE, False),
         ('Events:   Public reads published+expired | Managers read all | Managers create | ContentManagers update any', 8, LIGHT_GRAY, False),
         ('         ProgramManagers update own | Admins delete', 8, LIGHT_GRAY, False),
         ('', 3, WHITE, False),
         ('8 helper functions: user_has_role, user_is_admin, user_is_staff, user_has_min_role, auto_expire, auto_publish, ...', 8, MED_GRAY, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # ===== SLIDE 12: SECURITY =====
    slide = add_content_slide(prs, 'Security Architecture \u2014 Defense in Depth')

    layers = [
        ('Layer 1: Content Security Policy', 'Only approved scripts/styles/connections allowed \u2014 prevents XSS'),
        ('Layer 2: Row-Level Security (RLS)', 'Database enforces access on EVERY query \u2014 cannot be bypassed from client'),
        ('Layer 3: Client-Side Auth Guards', 'requireAuth() blocks pages, applyAuthVisibility() hides restricted UI'),
        ('Layer 4: Rate Limiting', '5 auth attempts per action per 60 seconds \u2014 prevents brute force'),
        ('Layer 5: Credential Isolation', 'config.js is .gitignored \u2014 never committed to repository'),
    ]

    for i, (title, desc) in enumerate(layers):
        y = Inches(1.1) + Inches(i * 0.78)
        colors = [RED, GREEN, TEAL, AMBER, MED_GRAY]
        add_shape_with_text(slide, Inches(0.5), y, Inches(9), Inches(0.65),
            [(title, 11, colors[i], True),
             (desc, 9, LIGHT_GRAY, False)],
            border_color=colors[i])

    # Key insight
    add_shape_with_text(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4),
        'Client-side checks = UX  |  Server-side RLS = Real security  |  DevTools cannot bypass database policies',
        font_size=9, fill_color=RGBColor(0x33, 0x2B, 0x00), font_color=GOLD,
        bold=True, alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 13: DEV WORKFLOW =====
    slide = add_content_slide(prs, 'Developer Workflow')

    add_shape_with_text(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(2.0),
        [('Getting Started', 13, GOLD, True),
         ('', 4, WHITE, False),
         ('1. git clone <repo>', 10, WHITE, False),
         ('2. cp config.example.js config.js', 10, WHITE, False),
         ('3. Edit config.js with Supabase keys', 10, WHITE, False),
         ('4. make serve', 10, TEAL, False),
         ('5. Open http://localhost:8000', 10, WHITE, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    add_shape_with_text(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(2.0),
        [('Build Commands', 13, GOLD, True),
         ('', 4, WHITE, False),
         ('make serve       \u2192  Dev server', 10, WHITE, False),
         ('make optimize    \u2192  WebP images', 10, WHITE, False),
         ('make minify      \u2192  Minify CSS/JS', 10, WHITE, False),
         ('make build       \u2192  Full prod build', 10, TEAL, False),
         ('make clean       \u2192  Remove artifacts', 10, WHITE, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # File map
    add_shape_with_text(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(1.8),
        [('Where to Look', 12, GOLD, True),
         ('', 3, WHITE, False),
         ('Page layout \u2192 *.html  |  Styles \u2192 styles.css  |  Site behavior \u2192 script.js', 9, LIGHT_GRAY, False),
         ('Auth logic \u2192 auth.js  |  Event CRUD \u2192 events.js  |  Admin UI \u2192 admin-events.js', 9, LIGHT_GRAY, False),
         ('Nav/Footer \u2192 common/  |  Team/Vendor data \u2192 data/*.json  |  DB schema \u2192 supabase-schema.sql', 9, LIGHT_GRAY, False),
         ('Supabase keys \u2192 config.js (never commit!)  |  Images \u2192 assets/', 9, LIGHT_GRAY, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # ===== SLIDE 14: TEST ACCOUNTS =====
    slide = add_content_slide(prs, 'Test Accounts')
    add_table(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(3.5),
        ['Email', 'Role', 'Password'],
        [
            ['rjramarao@gmail.com', 'GlobalAdmin', 'Menifee@92585'],
            ['orgadmin@aikyamusa.org', 'OrgAdmin', 'Menifee@92585'],
            ['contentmanager@aikyamusa.org', 'ContentManager', 'Menifee@92585'],
            ['programmanager@aikyamusa.org', 'ProgramManager', 'Menifee@92585'],
            ['financestaff@aikyamusa.org', 'FinanceStaff', 'Menifee@92585'],
            ['volunteer@aikyamusa.org', 'Volunteer', 'Menifee@92585'],
            ['donor@aikyamusa.org', 'Donor', 'Menifee@92585'],
            ['support@aikyamusa.org', 'Support', 'Menifee@92585'],
            ['member@aikyamusa.org', 'Member', 'Menifee@92585'],
        ],
        col_widths=[4.0, 2.5, 2.5])

    add_shape_with_text(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.4),
        'Login with different accounts to see how the UI adapts to each role!',
        font_size=11, fill_color=RGBColor(0x16, 0x20, 0x33), font_color=GOLD,
        bold=True, alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 15: WHAT'S NEXT =====
    slide = add_content_slide(prs, "What's Next")

    add_shape_with_text(slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(3.0),
        [('Immediate Roadmap', 13, GOLD, True),
         ('', 4, WHITE, False),
         ('\u2022  Event Registration / RSVP', 11, WHITE, False),
         ('\u2022  Gallery admin page', 11, WHITE, False),
         ('\u2022  Vendor admin page', 11, WHITE, False),
         ('\u2022  Member directory', 11, WHITE, False),
         ('\u2022  Email notifications', 11, WHITE, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    add_shape_with_text(slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(3.0),
        [('Future Possibilities', 13, GOLD, True),
         ('', 4, WHITE, False),
         ('\u2022  Payment integration (Stripe)', 11, WHITE, False),
         ('\u2022  PWA offline support', 11, WHITE, False),
         ('\u2022  Analytics dashboard', 11, WHITE, False),
         ('\u2022  Multi-language (i18n)', 11, WHITE, False),
         ('\u2022  Supabase Edge Functions', 11, WHITE, False)],
        border_color=RGBColor(0x33, 0x44, 0x66))

    # Scaling note
    add_shape_with_text(slide, Inches(0.5), Inches(4.4), Inches(9), Inches(0.8),
        [('Supabase Free Tier Limits (plenty of headroom)', 10, GOLD, True),
         ('50K auth users  |  500MB database  |  1GB storage  |  Unlimited API requests', 9, LIGHT_GRAY, False)],
        alignment=PP_ALIGN.CENTER, border_color=RGBColor(0x33, 0x44, 0x66))

    # ===== SLIDE 16: KEY TAKEAWAYS =====
    slide = add_content_slide(prs, 'Key Takeaways')

    takeaways = [
        ('ZERO COST', 'No servers, no monthly bills \u2014 static hosting + Supabase free tier'),
        ('SECURE', '9-role RBAC with database-level enforcement via Row-Level Security'),
        ('RESILIENT', 'JSON fallback ensures the site never shows a blank page'),
        ('SIMPLE', 'No frameworks, no build tools required \u2014 just HTML/CSS/JS'),
        ('EXTENSIBLE', 'New features = new JS module + DB table + RLS policies'),
    ]

    for i, (title, desc) in enumerate(takeaways):
        y = Inches(1.1) + Inches(i * 0.78)
        # Number
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.6), y + Inches(0.08), Inches(0.4), Inches(0.4))
        num.fill.solid()
        num.fill.fore_color.rgb = GOLD
        num.line.fill.background()
        tf = num.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK_BG
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_shape_with_text(slide, Inches(1.2), y, Inches(8.3), Inches(0.55),
            [(title, 13, GOLD, True),
             (desc, 10, LIGHT_GRAY, False)])

    # ===== SLIDE 17: THANK YOU =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0), Inches(1.8), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'Thank You'
    p.font.size = Pt(44)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = 'Questions & Discussion'
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(0), Inches(4.0), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = 'AIKYAM  |  Community Connecting Communities'
    p3.font.size = Pt(12)
    p3.font.color.rgb = MED_GRAY
    p3.alignment = PP_ALIGN.CENTER

    # ===== SAVE =====
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, '..', 'docs', 'AIKYAM_Architecture.pptx')
    out_path = os.path.abspath(out_path)
    prs.save(out_path)
    print(f'Saved: {out_path}')

if __name__ == '__main__':
    main()
